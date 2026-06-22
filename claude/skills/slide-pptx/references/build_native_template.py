"""
NATIVE PPTX TEMPLATE — matches the /slide HTML design system.
=============================================================
Editable text boxes, shapes, and tables. Substitutes:
  Manrope    -> Poppins   (install via: brew install --cask font-poppins)
  Public Sans -> Calibri
  JetBrains Mono -> Consolas
  OKLCH      -> sRGB approximations

Sections:
  1. Color + font tokens               (lines ~30-65)   adjust per brand
  2. Presentation setup                (lines ~67-78)   16:9 dimensions
  3. Helpers                           (lines ~80-200)  add_text, add_rect, add_block, ...
  4. Strip + foot                      (lines ~80-115)  slide # bar + page number
  5. Eyebrow + title + kicker          (lines ~115-145)
  6. Block variants                    (lines ~145-180) default/warm/dark/terra
  7. Bullet + numbered list            (lines ~180-220)
  8. Table helper                      (lines ~220-280)
  9. Six worked patterns (slides 1-10): cover, cadre, comparatif, wide table,
     flow+sidebar, action plan. Each pattern is a self-contained section
     between `# SLIDE N` markers; copy the ones you need, replace content,
     delete the rest.

To adapt for a new deck:
  - Copy this file to your working directory
  - Edit OUTPUT path at the bottom
  - For each slide, keep the pattern that fits and replace the content
  - Remove unused slides

Run: /tmp/pptx-venv/bin/python build_native.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# === COLOR TOKENS ===
PAPER       = RGBColor(0xFD, 0xFD, 0xFD)
PAPER_SOFT  = RGBColor(0xF6, 0xF7, 0xFA)
PAPER_WARM  = RGBColor(0xF7, 0xF2, 0xE2)
PAPER_DEEP  = RGBColor(0xEC, 0xEE, 0xF1)
NAVY        = RGBColor(0x1A, 0x1F, 0x3A)
INK         = RGBColor(0x1A, 0x1F, 0x3A)
INK_BODY    = RGBColor(0x35, 0x38, 0x4B)
INK_SOFT    = RGBColor(0x60, 0x63, 0x76)
INK_MUTE    = RGBColor(0x8B, 0x8D, 0x99)
RULE        = RGBColor(0xD8, 0xD9, 0xDD)
RULE_STRONG = RGBColor(0xB6, 0xB8, 0xBE)
YELLOW      = RGBColor(0xEB, 0xC8, 0x45)
YELLOW_DEEP = RGBColor(0xA8, 0x78, 0x10)
YELLOW_SOFT = RGBColor(0xF6, 0xEC, 0xC8)
TERRA       = RGBColor(0x96, 0x4A, 0x2E)
TERRA_SOFT  = RGBColor(0xEF, 0xDB, 0xD0)
GREEN       = RGBColor(0x4D, 0x76, 0x5E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

DISPLAY = 'Poppins'   # Manrope substitute
BODY    = 'Calibri'   # Public Sans substitute
MONO    = 'Consolas'  # JetBrains Mono substitute

# === PRESENTATION SETUP ===
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

BLANK = prs.slide_layouts[6]

# === UNITS ===
def IN(x): return Inches(x)

# === HELPERS ===
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x), IN(y), IN(w), IN(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s

def add_line_h(slide, x, y, w, color, weight=0.5):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x), IN(y), IN(w), Pt(weight))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, runs, align='left', anchor='top', vmargin=0.04):
    """runs: list of (text, dict opts). opts keys: font, size, bold, color, italic, spacing, kern."""
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(vmargin * 72)
    tf.margin_bottom = Pt(vmargin * 72)
    if anchor == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    # First paragraph already exists
    if isinstance(runs, str):
        runs = [(runs, {})]
    # Group runs by paragraph (use \n as separator inside text)
    paragraphs = [[]]
    for text, opts in runs:
        if text is None:
            paragraphs.append([])
            continue
        paragraphs[-1].append((text, opts))

    for pi, parts in enumerate(paragraphs):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if align == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        for text, opts in parts:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = opts.get('font', BODY)
            f.size = Pt(opts.get('size', 11))
            f.bold = opts.get('bold', False)
            f.italic = opts.get('italic', False)
            if 'color' in opts:
                f.color.rgb = opts['color']
        if 'space_after' in (parts[0][1] if parts else {}):
            p.space_after = Pt(parts[0][1]['space_after'])
    return tb

# ===================== STRIP / FOOT =====================
STRIP_H = 0.42
FOOT_H  = 0.32

def add_strip(slide, left, topic, brand_main='Theodo', brand_sub='HealthTech'):
    add_rect(slide, 0, 0, 13.333, STRIP_H, PAPER_SOFT)
    add_line_h(slide, 0, STRIP_H, 13.333, RULE, 0.5)
    add_text(slide, 0.4, 0, 6, STRIP_H,
             [(left, {'font': MONO, 'size': 9, 'color': INK_MUTE}),
              ('  / ', {'font': MONO, 'size': 9, 'color': INK_MUTE}),
              (topic, {'font': MONO, 'size': 9, 'color': YELLOW_DEEP, 'bold': True})],
             anchor='middle')
    add_text(slide, 9, 0, 4, STRIP_H,
             [(brand_main, {'font': DISPLAY, 'size': 10, 'color': NAVY, 'bold': True}),
              (' • ', {'font': DISPLAY, 'size': 10, 'color': YELLOW_DEEP, 'bold': True}),
              (brand_sub, {'font': DISPLAY, 'size': 10, 'color': NAVY, 'bold': True})],
             align='right', anchor='middle')

def add_foot(slide, num_str, mid, right):
    y = 7.5 - FOOT_H
    add_rect(slide, 0, y, 13.333, FOOT_H, PAPER_SOFT)
    add_line_h(slide, 0, y, 13.333, RULE, 0.5)
    add_text(slide, 0.4, y, 2, FOOT_H,
             [(num_str, {'font': MONO, 'size': 8, 'color': NAVY, 'bold': True})],
             anchor='middle')
    add_text(slide, 2.5, y, 7, FOOT_H,
             [(mid, {'font': MONO, 'size': 8, 'color': INK_MUTE})],
             anchor='middle')
    add_text(slide, 9, y, 4, FOOT_H,
             [(right, {'font': MONO, 'size': 8, 'color': INK_MUTE})],
             align='right', anchor='middle')

# ===================== BLOCK / EYEBROW / TITLE / KICKER =====================
def add_eyebrow(slide, x, y, text, w=6):
    add_text(slide, x, y, w, 0.32,
             [(text.upper(), {'font': MONO, 'size': 9, 'color': YELLOW_DEEP, 'bold': True})])

def add_title(slide, x, y, w, h, parts):
    """parts: list of (text, color, bold). Color None = NAVY."""
    runs = []
    for text, color in parts:
        runs.append((text, {'font': DISPLAY, 'size': 30, 'color': color or NAVY, 'bold': True}))
    add_text(slide, x, y, w, h, runs)

def add_kicker(slide, x, y, w, h, runs_input):
    """runs_input: string or list of (text, opts) tuples."""
    if isinstance(runs_input, str):
        runs = [(runs_input, {'font': BODY, 'size': 14, 'color': INK_SOFT})]
    else:
        runs = []
        for item in runs_input:
            if isinstance(item, str):
                runs.append((item, {'font': BODY, 'size': 14, 'color': INK_SOFT}))
            else:
                text, bold = item
                runs.append((text, {'font': BODY, 'size': 14, 'color': NAVY if bold else INK_SOFT, 'bold': bold}))
    add_text(slide, x, y, w, h, runs)

def add_block(slide, x, y, w, h, kind='default'):
    """kind: default, warm, dark, terra, green."""
    if kind == 'dark':
        fill = NAVY
        return add_rect(slide, x, y, w, h, fill)
    elif kind == 'warm':
        fill = PAPER_WARM
        s = add_rect(slide, x, y, w, h, fill, line=YELLOW)
        return s
    elif kind == 'terra':
        fill = TERRA_SOFT
        s = add_rect(slide, x, y, w, h, fill, line=TERRA)
        return s
    else:
        fill = PAPER_SOFT
        s = add_rect(slide, x, y, w, h, fill, line=RULE)
        return s

def block_header(slide, x, y, w, text, color=YELLOW_DEEP):
    add_text(slide, x, y, w, 0.3,
             [(text.upper(), {'font': MONO, 'size': 9, 'color': color, 'bold': True})])

def block_dark_header(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.3,
             [(text.upper(), {'font': MONO, 'size': 9, 'color': YELLOW, 'bold': True})])

# bullet helper
def bullet_list(slide, x, y, w, h, items, dark=False, size=11):
    """items: list of strings (or list of (str, bool_strong) tuples for partial bold via simple syntax)."""
    runs = []
    for i, item in enumerate(items):
        if i > 0:
            runs.append((None, {}))  # paragraph break
        runs.append(('•  ', {'font': BODY, 'size': size, 'color': YELLOW if dark else NAVY, 'bold': True}))
        if isinstance(item, str):
            runs.append((item, {'font': BODY, 'size': size, 'color': WHITE if dark else INK_BODY}))
        else:
            for chunk in item:
                if isinstance(chunk, str):
                    runs.append((chunk, {'font': BODY, 'size': size, 'color': WHITE if dark else INK_BODY}))
                else:
                    text, bold = chunk
                    runs.append((text, {'font': BODY, 'size': size, 'color': WHITE if dark else INK, 'bold': bold}))
    add_text(slide, x, y, w, h, runs)

def numbered_list(slide, x, y, w, h, items, size=11):
    runs = []
    for i, item in enumerate(items):
        if i > 0:
            runs.append((None, {}))
        runs.append(('%02d  ' % (i+1), {'font': MONO, 'size': size-1, 'color': YELLOW_DEEP, 'bold': True}))
        if isinstance(item, str):
            runs.append((item, {'font': BODY, 'size': size, 'color': INK_BODY}))
        else:
            for chunk in item:
                if isinstance(chunk, str):
                    runs.append((chunk, {'font': BODY, 'size': size, 'color': INK_BODY}))
                else:
                    text, bold = chunk
                    runs.append((text, {'font': BODY, 'size': size, 'color': INK, 'bold': bold}))
    add_text(slide, x, y, w, h, runs)

# ===================== TABLE =====================
def add_table(slide, x, y, w, h, rows, col_widths=None, header=None, hilite_rows=None):
    """rows: list of lists. header: list of column names (optional)."""
    nrows = len(rows) + (1 if header else 0)
    ncols = len(rows[0]) if rows else 1
    tbl_shape = slide.shapes.add_table(nrows, ncols, IN(x), IN(y), IN(w), IN(h))
    tbl = tbl_shape.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = IN(cw)
    hilite_rows = hilite_rows or []
    r_offset = 0
    if header:
        for ci, hdr in enumerate(header):
            cell = tbl.cell(0, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER_SOFT
            tf = cell.text_frame
            tf.margin_left = Pt(4); tf.margin_right = Pt(4)
            tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
            tf.text = ''
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = hdr.upper()
            r.font.name = MONO
            r.font.size = Pt(8)
            r.font.bold = True
            r.font.color.rgb = INK_MUTE
        r_offset = 1
    for ri, row in enumerate(rows):
        is_hilite = ri in hilite_rows
        for ci, cell_data in enumerate(row):
            cell = tbl.cell(ri + r_offset, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = YELLOW_SOFT if is_hilite else PAPER
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(4); tf.margin_right = Pt(4)
            tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
            tf.text = ''
            p = tf.paragraphs[0]
            # cell_data: string or list of (text, opts)
            if isinstance(cell_data, str):
                cell_data = [(cell_data, {})]
            for text, opts in cell_data:
                r = p.add_run()
                r.text = text
                r.font.name = opts.get('font', BODY)
                r.font.size = Pt(opts.get('size', 10))
                r.font.bold = opts.get('bold', False)
                r.font.italic = opts.get('italic', False)
                if 'color' in opts:
                    r.font.color.rgb = opts['color']
    return tbl_shape

# ============================================================
# SLIDE 1 · COVER
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.333, 7.5, PAPER)
add_strip(s, 'BRIEFING TECHNIQUE', 'IDENTIFICATION ÉLECTRONIQUE')
add_foot(s, '01 / 10', '', 'BRIEFING · IEU & IEPS')

# Meta line top
add_text(s, 1.2, 0.85, 8, 0.3,
         [('SUNRISE', {'font': MONO, 'size': 10, 'color': INK_MUTE}),
          ('  ·  ', {'font': MONO, 'size': 10, 'color': YELLOW_DEEP}),
          ('DMN V1.2.2', {'font': MONO, 'size': 10, 'color': INK_MUTE}),
          ('  ·  ', {'font': MONO, 'size': 10, 'color': YELLOW_DEEP}),
          ('26 MAI 2026', {'font': MONO, 'size': 10, 'color': INK_MUTE})])

# Eyebrow
add_text(s, 1.2, 2.6, 6, 0.3,
         [('BRIEFING · 10 SLIDES', {'font': MONO, 'size': 10, 'color': YELLOW_DEEP, 'bold': True})])

# Big title
add_text(s, 1.2, 2.95, 11, 1.8,
         [('IEU ', {'font': DISPLAY, 'size': 92, 'color': NAVY, 'bold': True}),
          ('& ', {'font': DISPLAY, 'size': 92, 'color': TERRA}),
          ('IEPS', {'font': DISPLAY, 'size': 92, 'color': YELLOW_DEEP, 'bold': True})])

# Lead
add_text(s, 1.2, 4.95, 11, 0.8,
         [("Identifier les utilisateurs d'un DMN selon le PGSSI‑S.",
           {'font': DISPLAY, 'size': 22, 'color': INK, 'bold': False})])

# Author
add_text(s, 1.2, 5.65, 11, 0.4,
         [('Pour Sunrise · Theodo HealthTech QARA',
           {'font': BODY, 'size': 13, 'color': INK_SOFT, 'italic': True})])

# Cover foot
add_text(s, 1.2, 6.8, 8, 0.3,
         [('Exigences_référentiel_FR_DMN_V1.2.2 · ANS · 2023‑02‑22',
           {'font': MONO, 'size': 10, 'color': INK_MUTE})])
add_text(s, 9.5, 6.7, 3.5, 0.4,
         [('Sunrise ', {'font': DISPLAY, 'size': 18, 'color': NAVY, 'bold': True}),
          ('× ', {'font': DISPLAY, 'size': 18, 'color': YELLOW_DEEP}),
          ('Theodo', {'font': DISPLAY, 'size': 18, 'color': NAVY, 'bold': True})],
         align='right')

# ============================================================
# SLIDE 2 · CADRE
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.333, 7.5, PAPER)
add_strip(s, '02 · CADRE', 'IEU & IEPS DANS LA GRILLE DMN')
add_foot(s, '02 / 10', 'CADRE · RÉFÉRENTIEL DMN V1.2.2', 'PGSSI‑S · ANS')

add_eyebrow(s, 0.5, 0.7, 'Cadre')
add_text(s, 0.5, 1.0, 12, 0.8,
         [('IEU & ', {'font': DISPLAY, 'size': 30, 'color': NAVY, 'bold': True}),
          ('IEPS', {'font': DISPLAY, 'size': 30, 'color': YELLOW_DEEP, 'bold': True}),
          (' dans la grille DMN ', {'font': DISPLAY, 'size': 30, 'color': NAVY, 'bold': True}),
          ('v1.2.2', {'font': DISPLAY, 'size': 30, 'color': YELLOW_DEEP, 'bold': True})])
add_kicker(s, 0.5, 1.65, 12, 0.5,
           [('84 exigences. 8 profils. 9 sections. IEU et IEPS couvrent l\'identification de tout utilisateur du système.', False)])

# Left: dark block "La règle DMN"
add_block(s, 0.5, 2.5, 6, 4.3, 'dark')
block_dark_header(s, 0.75, 2.65, 5.5, 'La règle DMN')
add_text(s, 0.75, 3.0, 5.5, 0.4,
         [('Pour chaque utilisateur, démontrer :', {'font': BODY, 'size': 12, 'color': WHITE})])
bullet_list(s, 0.75, 3.45, 5.5, 3.0,
            [[('Identification ', False), ('fiable, traçable, unique', True)],
             [('Authentification au niveau ', False), ('PGSSI‑S', True), (' requis', False)],
             [('Cycle de vie MIE : création, renouvellement, révocation, déconnexion', False)]],
            dark=True, size=12)

# Right: IEPS block + IEU block + note
add_block(s, 6.85, 2.5, 6.0, 1.55, 'default')
block_header(s, 7.1, 2.65, 5.5, 'IEPS · 9 EXIGENCES')
add_text(s, 7.1, 3.0, 5.5, 0.9,
         [('Professionnels de santé. Socle : ', {'font': BODY, 'size': 12, 'color': INK_BODY}),
          ('PSC + Annuaire Santé + PGSSI‑S IE ASPP', {'font': BODY, 'size': 12, 'color': INK, 'bold': True}),
          ('.', {'font': BODY, 'size': 12, 'color': INK_BODY})])

add_block(s, 6.85, 4.25, 6.0, 1.55, 'warm')
block_header(s, 7.1, 4.4, 5.5, 'IEU · 12 EXIGENCES')
add_text(s, 7.1, 4.75, 5.5, 0.9,
         [('Patients / usagers. Socle : ', {'font': BODY, 'size': 12, 'color': INK_BODY}),
          ('INS + 2FA + PGSSI‑S IE Usagers', {'font': BODY, 'size': 12, 'color': INK, 'bold': True}),
          ('.', {'font': BODY, 'size': 12, 'color': INK_BODY})])

add_text(s, 6.85, 6.0, 6.0, 0.5,
         [('Sunrise adresse PS + patients : les deux chapitres s\'appliquent.',
           {'font': BODY, 'size': 11, 'color': INK_SOFT, 'italic': True})])

# ============================================================
# SLIDE 3 · COMPARATIF
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.333, 7.5, PAPER)
add_strip(s, '03 · COMPARATIF', 'IEPS VS IEU')
add_foot(s, '03 / 10', 'COMPARATIF · IEPS VS IEU', 'DMN V1.2.2')

add_eyebrow(s, 0.5, 0.7, 'Comparatif')
add_text(s, 0.5, 1.0, 12, 0.8,
         [('IEPS et IEU : ', {'font': DISPLAY, 'size': 30, 'color': NAVY, 'bold': True}),
          ('tableau de bord', {'font': DISPLAY, 'size': 30, 'color': YELLOW_DEEP, 'bold': True})])

# Left card: IEPS
add_block(s, 0.5, 2.1, 6.15, 4.6, 'default')
# Left border navy
add_rect(s, 0.5, 2.1, 0.06, 4.6, NAVY)
add_text(s, 0.75, 2.25, 5.5, 0.5,
         [('IEPS · Professionnels', {'font': DISPLAY, 'size': 18, 'color': NAVY, 'bold': True})])
ips_rows = [
    ('Profil', 'Accès Professionnel'),
    ('Exigences', '9 (IEPS 2, 4-9, 12, 13)'),
    ('Identifiant', 'RPPS / ADELI (Annuaire Santé)'),
    ('Auth.', 'PSC, acr_values=eidas1'),
    ('Socle', 'PGSSI‑S IE ASPP'),
    ('eIDAS', 'Substantiel (PSC OTP) / Élevé (CPx)'),
]
y0 = 2.85
for k, v in ips_rows:
    add_text(s, 0.75, y0, 1.6, 0.35,
             [(k, {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})], anchor='middle')
    add_text(s, 2.4, y0, 4.0, 0.35,
             [(v, {'font': BODY, 'size': 11, 'color': INK_BODY})], anchor='middle')
    add_line_h(s, 0.75, y0 + 0.35, 5.5, RULE, 0.4)
    y0 += 0.5

# Right card: IEU
add_block(s, 6.85, 2.1, 6.0, 4.6, 'warm')
add_rect(s, 6.85, 2.1, 0.06, 4.6, YELLOW_DEEP)
add_text(s, 7.1, 2.25, 5.5, 0.5,
         [('IEU · Usagers', {'font': DISPLAY, 'size': 18, 'color': YELLOW_DEEP, 'bold': True})])
ieu_rows = [
    ('Profil', 'Accès Usager (+ ApCV opt.)'),
    ('Exigences', '12 (IEU 1 à 12)'),
    ('Identifiant', 'Matricule INS (IEU 7)'),
    ('Auth.', '2FA obligatoire (IEU 9)'),
    ('Socle', 'PGSSI‑S IE Usagers + Guide INS V3.0'),
    ('eIDAS', 'Substantiel mini (FC+, ApCV, OTP)'),
]
y0 = 2.85
for k, v in ieu_rows:
    add_text(s, 7.1, y0, 1.6, 0.35,
             [(k, {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})], anchor='middle')
    add_text(s, 8.75, y0, 4.0, 0.35,
             [(v, {'font': BODY, 'size': 11, 'color': INK_BODY})], anchor='middle')
    add_line_h(s, 7.1, y0 + 0.35, 5.5, RULE, 0.4)
    y0 += 0.5

# Note bottom
add_text(s, 0.5, 6.85, 12.3, 0.3,
         [('IEPS s\'appuie sur PSC. IEU s\'appuie sur INS et impose la 2FA. Les deux héritent du PGSSI‑S.',
           {'font': BODY, 'size': 11, 'color': INK_SOFT, 'italic': True})],
         align='center')

# ============================================================
# SLIDE 4 · IEPS CADRE
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.333, 7.5, PAPER)
add_strip(s, '04 · IEPS', 'CADRE')
add_foot(s, '04 / 10', 'IEPS · CADRE', 'PSC · ANNUAIRE SANTÉ')

add_eyebrow(s, 0.5, 0.7, 'IEPS · Cadre')
add_text(s, 0.5, 1.0, 12, 0.8,
         [('IEPS : identification des ', {'font': DISPLAY, 'size': 30, 'color': NAVY, 'bold': True}),
          ('professionnels', {'font': DISPLAY, 'size': 30, 'color': YELLOW_DEEP, 'bold': True})])
add_kicker(s, 0.5, 1.65, 12, 0.6,
           'Tout PS connecté à Sunrise déclenche les 9 exigences IEPS. Authentification de référence : Pro Santé Connect.')

# Left column: Quand + Authentification
add_text(s, 0.5, 2.5, 6, 0.35,
         [('Quand', {'font': DISPLAY, 'size': 16, 'color': INK, 'bold': True})])
bullet_list(s, 0.5, 2.85, 6, 1.0,
            ["À chaque connexion d'un PS : médecin, IDE, secrétaire médicale, kiné, autre ASPP.",
             "Y compris comptes back-office et SSO d'établissement."],
            size=11)
add_text(s, 0.5, 4.0, 6, 0.35,
         [('Authentification attendue', {'font': DISPLAY, 'size': 16, 'color': INK, 'bold': True})])
bullet_list(s, 0.5, 4.35, 6, 2.4,
            [[('Identifiant', True), (' : RPPS / ADELI (Annuaire Santé, ANN 1-5).', False)],
             [('PSC', True), (' : mode web, native+navigateur, ou CIBA (PSC 2).', False)],
             [('acr_values', True), (' = eidas1 (PSC 6).', False)],
             [('Eligibilité + engagement + CGU', True), (' : PSC 1, 3, 5.', False)]],
            size=11)

# Right column: dark + terra
add_block(s, 6.85, 2.5, 6.0, 2.0, 'dark')
block_dark_header(s, 7.1, 2.65, 5.5, 'Référentiels socles')
bullet_list(s, 7.1, 3.0, 5.5, 1.5,
            ['PGSSI‑S IE ASPP',
             'Référentiel Pro Santé Connect',
             'Annuaire Santé (MAJ quotidienne)'],
            dark=True, size=11)

add_block(s, 6.85, 4.7, 6.0, 2.0, 'terra')
block_header(s, 7.1, 4.85, 5.5, 'Pièges fréquents', color=TERRA)
bullet_list(s, 7.1, 5.2, 5.5, 1.5,
            [[('Compte legacy (email/MdP) : à rapprocher à la 1ʳᵉ connexion PSC (PSC 4). Sinon doublon = ', False), ('IEPS 8 KO', True), ('.', False)],
             ['Unicité email seule = limite. Ajouter nom + prénom dans la clé d\'unicité.']],
            size=10)

# ============================================================
# SLIDE 5 · IEPS TABLE
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.333, 7.5, PAPER)
add_strip(s, '05 · IEPS', 'LES 9 EXIGENCES')
add_foot(s, '05 / 10', 'IEPS · 9 EXIGENCES', 'DMN V1.2.2')

add_eyebrow(s, 0.5, 0.7, 'IEPS · Exigences')
add_text(s, 0.5, 1.0, 12, 0.7,
         [('Les 9 exigences ', {'font': DISPLAY, 'size': 30, 'color': NAVY, 'bold': True}),
          ('IEPS', {'font': DISPLAY, 'size': 30, 'color': YELLOW_DEEP, 'bold': True})])
add_kicker(s, 0.5, 1.55, 12, 0.5,
           'Trois blocs : gestion du compte, identification, connexion / déconnexion.')

# Table 9 rows
rows = [
    [[('IEPS 2', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Email ', {'size': 10}), ('ou ', {'size': 10, 'bold': True}), ('tél vérifié à création (l\'un suffit).', {'size': 10})],
     [('Double opt-in à création seulement. Compte non verrouillé si rappel envoyé.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEPS 4', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Attributs minimum : nom d\'exercice, prénom d\'exercice (Annuaire), profession.', {'size': 10})],
     [('Modèle aligné Annuaire Santé. Profils existants : ', {'size': 9, 'color': INK_SOFT}),
      ('campagne de vérif rétroactive', {'size': 9, 'color': INK, 'bold': True}),
      ('.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEPS 5', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Renouvellement MdP à la demande.', {'size': 10})],
     [('Lien « MdP oublié », jeton à durée limitée.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEPS 6', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Renouvellement MIE à fréquence paramétrable.', {'size': 10})],
     [('Pas d\'exigence stricte. Reco ANS : 6 mois. Politique admin configurable.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEPS 7', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Révocation par le PS de tous ses MIE.', {'size': 10})],
     [('Auto-révocation en cas de perte / vol / compromission. Inclut « déconnecter toutes les sessions ».', {'size': 9, 'color': INK_SOFT})]],
    [[('IEPS 8', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Unicité des identifiants.', {'size': 10})],
     [('Email seul = limite. Recommandé : email + nom + prénom.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEPS 9', {'font': MONO, 'size': 10, 'color': TERRA, 'bold': True})],
     [('Complexité MdP + restriction d\'accès (PGSSI‑S IE ASPP).', {'size': 10})],
     [('Anti-brute force (Sunrise : 5 tentatives / 5 min ✓). Entropie MdP bloquée à création. Preuve = vidéo.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEPS 12', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Déconnexion à la demande.', {'size': 10})],
     [('Invalidation session + refresh token côté serveur.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEPS 13', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Déconnexion auto après inactivité paramétrable.', {'size': 10})],
     [('Idle timeout front + serveur.', {'size': 9, 'color': INK_SOFT})]],
]
add_table(s, 0.5, 2.25, 12.35, 4.5,
          rows, col_widths=[1.05, 5.0, 6.3],
          header=['ID', 'Énoncé', 'Implication produit'],
          hilite_rows=[6])

# ============================================================
# SLIDE 6 · IEPS PILE
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.333, 7.5, PAPER)
add_strip(s, '06 · IEPS', 'PILE TECHNIQUE')
add_foot(s, '06 / 10', 'IEPS · PILE', 'PSC · ANNUAIRE · PGSSI‑S')

add_eyebrow(s, 0.5, 0.7, 'IEPS · Pile technique')
add_text(s, 0.5, 1.0, 12, 0.8,
         [('PSC pour l\'auth, Annuaire pour ', {'font': DISPLAY, 'size': 28, 'color': NAVY, 'bold': True}),
          ('l\'identité', {'font': DISPLAY, 'size': 28, 'color': YELLOW_DEEP, 'bold': True})])
add_kicker(s, 0.5, 1.7, 12, 0.6,
           'Authentification déléguée à PSC. Identité hydratée depuis l\'Annuaire Santé. Aucun MdP PS stocké côté Sunrise.')

# Left: flow diagram
add_text(s, 0.5, 2.5, 6, 0.35,
         [('Flux PSC (mode web)', {'font': DISPLAY, 'size': 16, 'color': INK, 'bold': True})])

# Flow box backdrop
add_block(s, 0.5, 2.95, 6.3, 4.0, 'default')

# Helper for flow nodes
def flow_node(x, y, w, h, text, kind='default'):
    if kind == 'k':
        fill, fg = NAVY, WHITE
    elif kind == 'y':
        fill, fg = YELLOW_SOFT, NAVY
    elif kind == 't':
        fill, fg = TERRA_SOFT, TERRA
    else:
        fill, fg = PAPER, INK_BODY
    add_rect(s, x, y, w, h, fill, line=RULE_STRONG)
    add_text(s, x, y, w, h,
             [(text, {'font': MONO, 'size': 9, 'color': fg, 'bold': kind in ('k', 'y')})],
             align='center', anchor='middle')

def flow_arr(x, y):
    add_text(s, x, y, 0.3, 0.3,
             [('→', {'font': BODY, 'size': 12, 'color': INK_MUTE, 'bold': True})], align='center', anchor='middle')

# Row 1: Sunrise → PSC /authorize → IdN PSC
flow_node(0.75, 3.15, 1.5, 0.35, 'Sunrise')
flow_arr(2.3, 3.15)
flow_node(2.65, 3.15, 1.8, 0.35, 'PSC /authorize', 'y')
flow_arr(4.5, 3.15)
flow_node(4.85, 3.15, 1.6, 0.35, 'IdN PSC', 'k')

# Row 2: PSC → CPx / eCPS / OTP / CIBA
flow_node(0.75, 3.85, 1.0, 0.35, 'PSC', 'k')
flow_arr(1.8, 3.85)
flow_node(2.15, 3.85, 3.8, 0.35, 'CPx / eCPS / OTP / CIBA')

# Row 3: PSC → Callback + ID Token → Sunrise
flow_node(0.75, 4.55, 1.0, 0.35, 'PSC', 'k')
flow_arr(1.8, 4.55)
flow_node(2.15, 4.55, 2.2, 0.35, 'Callback + ID Token', 'y')
flow_arr(4.4, 4.55)
flow_node(4.75, 4.55, 1.6, 0.35, 'Sunrise')

# Row 4: Sunrise → Annuaire Santé → Profil PS
flow_node(0.75, 5.25, 1.4, 0.35, 'Sunrise')
flow_arr(2.2, 5.25)
flow_node(2.55, 5.25, 2.0, 0.35, 'Annuaire Santé', 't')
flow_arr(4.6, 5.25)
flow_node(4.95, 5.25, 1.6, 0.35, 'Profil PS', 'y')

# Right: Annuaire block + Watch-outs block
add_block(s, 7.0, 2.95, 5.85, 1.85, 'default')
block_header(s, 7.25, 3.1, 5.4, 'Annuaire Santé · ANN 1-5')
bullet_list(s, 7.25, 3.45, 5.4, 1.4,
            ['Récupérer RPPS + adresse MSSanté.',
             'Rafraîchir adossement (MAJ J‑1 mois max).',
             'Associer correspondants + utilisateurs.'],
            size=10)

add_block(s, 7.0, 5.0, 5.85, 1.95, 'terra')
block_header(s, 7.25, 5.15, 5.4, 'Watch-outs', color=TERRA)
bullet_list(s, 7.25, 5.5, 5.4, 1.4,
            [[('acr_values', True), (' ≠ ', False), ('eidas1', True), (' = rejet.', False)],
             ['CGU PSC absentes = PSC 5 KO.'],
             ['Compte legacy non rapproché = PSC 4 + IEPS 8 KO.']],
            size=10)

# ============================================================
# SLIDE 7 · IEU CADRE
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.333, 7.5, PAPER)
add_strip(s, '07 · IEU', 'CADRE')
add_foot(s, '07 / 10', 'IEU · CADRE', 'INS · PGSSI‑S IE USAGERS')

add_eyebrow(s, 0.5, 0.7, 'IEU · Cadre')
add_text(s, 0.5, 1.0, 12, 0.8,
         [('IEU : identification des ', {'font': DISPLAY, 'size': 30, 'color': NAVY, 'bold': True}),
          ('usagers', {'font': DISPLAY, 'size': 30, 'color': YELLOW_DEEP, 'bold': True})])
add_kicker(s, 0.5, 1.65, 12, 0.6,
           'Tout patient connecté à Sunrise déclenche les 12 exigences IEU. Pivot : matricule INS. 2FA obligatoire.')

add_text(s, 0.5, 2.5, 6, 0.35,
         [('Quand', {'font': DISPLAY, 'size': 16, 'color': INK, 'bold': True})])
bullet_list(s, 0.5, 2.85, 6, 1.0,
            ['Portail patient, app de télésuivi, saisie de symptômes.',
             "Auto-inscription ou compte créé par le médecin et activé par l'usager."],
            size=11)
add_text(s, 0.5, 4.0, 6, 0.35,
         [('Authentification attendue', {'font': DISPLAY, 'size': 16, 'color': INK, 'bold': True})])
bullet_list(s, 0.5, 4.35, 6, 2.4,
            [[('Identifiant', True), (' : matricule INS (IEU 7).', False)],
             [('2 facteurs', True), (' imposés (IEU 9).', False)],
             [('eIDAS', True), (' substantiel minimum.', False)],
             ['Lien fort avec section INS (qualification, traçabilité).']],
            size=11)

# Right: warm block (INS pivot) + dark block (2FA)
add_block(s, 6.85, 2.5, 6.0, 2.5, 'warm')
block_header(s, 7.1, 2.65, 5.5, 'Identifiant pivot : matricule INS')
add_text(s, 7.1, 3.0, 5.5, 1.0,
         [('IEU 7', {'font': BODY, 'size': 11, 'color': INK, 'bold': True}),
          (' : supporter le matricule INS. ', {'font': BODY, 'size': 11, 'color': INK_BODY}),
          ('IEU 8', {'font': BODY, 'size': 11, 'color': INK, 'bold': True}),
          (" : aller le chercher si l'IdN ne le donne pas (INSi ou flux IHE PAM).",
           {'font': BODY, 'size': 11, 'color': INK_BODY})])
add_text(s, 7.1, 4.05, 5.5, 0.7,
         [("Si l'usager saisit ses traits, Sunrise devient référentiel d'identités (Guide INS V3.0 §1.2).",
           {'font': BODY, 'size': 10, 'color': INK_SOFT, 'italic': True})])

add_block(s, 6.85, 5.15, 6.0, 1.7, 'dark')
block_dark_header(s, 7.1, 5.3, 5.5, '2FA obligatoire (IEU 9)')
add_text(s, 7.1, 5.65, 5.5, 1.0,
         [('Pas de MdP seul. Valides : MdP+OTP, FranceConnect+, ApCV, eID eIDAS.',
           {'font': BODY, 'size': 11, 'color': WHITE})])

# ============================================================
# SLIDE 8 · IEU TABLE
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.333, 7.5, PAPER)
add_strip(s, '08 · IEU', 'LES 12 EXIGENCES')
add_foot(s, '08 / 10', 'IEU · 12 EXIGENCES', 'DMN V1.2.2')

add_eyebrow(s, 0.5, 0.7, 'IEU · Exigences')
add_text(s, 0.5, 1.0, 12, 0.7,
         [('Les 12 exigences ', {'font': DISPLAY, 'size': 30, 'color': NAVY, 'bold': True}),
          ('IEU', {'font': DISPLAY, 'size': 30, 'color': YELLOW_DEEP, 'bold': True})])
add_kicker(s, 0.5, 1.55, 12, 0.55,
           'Trois blocs : gestion du compte (1-6), identification (7-8), connexion / déconnexion (9-12). IEU 7 et IEU 9 = exigences clés.')

rows = [
    [[('IEU 1', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [("Vérification des attributs d'identité à création.", {'size': 10})],
     [('Justificatif demandé ou flux INS reçu.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEU 2', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [("Modification d'identité = vérification équivalente.", {'size': 10})],
     [("Procédure équivalente à l'inscription initiale.", {'size': 9, 'color': INK_SOFT})]],
    [[('IEU 3', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Email / tél vérifiés à création (si auth ou récup).', {'size': 10})],
     [('Double opt-in, OTP sur canal renseigné.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEU 4', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Renouvellement MdP à la demande.', {'size': 10})],
     [('Lien « MdP oublié », jeton à durée limitée.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEU 5', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Renouvellement MIE à fréquence paramétrable.', {'size': 10})],
     [('Politique admin configurable.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEU 6', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [("Révocation par l'usager d'un de ses MIE.", {'size': 10})],
     [('Désinscrire un téléphone OTP perdu.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEU 7', {'font': MONO, 'size': 10, 'color': TERRA, 'bold': True})],
     [('Matricule INS comme identifiant.', {'size': 10})],
     [('Stocker matricule INS + OID. Indexer le compte.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEU 8', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [("Recherche du matricule INS si IdN ne fournit qu'un ID privé.", {'size': 10})],
     [('Appel INSi ou intégration flux IHE PAM.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEU 9', {'font': MONO, 'size': 10, 'color': TERRA, 'bold': True})],
     [('2FA imposée.', {'size': 10})],
     [('MdP+OTP, FranceConnect+, ApCV, ou eID reconnu.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEU 10', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [("Complexité MdP + restriction d'accès (PGSSI‑S).", {'size': 10})],
     [('Anti-brute force, longueur / complexité, journal.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEU 11', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Déconnexion à la demande.', {'size': 10})],
     [('Invalidation côté serveur.', {'size': 9, 'color': INK_SOFT})]],
    [[('IEU 12', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Déconnexion auto après inactivité paramétrable.', {'size': 10})],
     [('Idle timeout front + serveur.', {'size': 9, 'color': INK_SOFT})]],
]
add_table(s, 0.5, 2.2, 12.35, 4.6,
          rows, col_widths=[1.05, 5.0, 6.3],
          header=['ID', 'Énoncé', 'Implication produit'],
          hilite_rows=[6, 8])

# ============================================================
# SLIDE 9 · IEU PILE
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.333, 7.5, PAPER)
add_strip(s, '09 · IEU', 'PILE TECHNIQUE')
add_foot(s, '09 / 10', 'IEU · PILE', '2FA · INSi · FC+ · APCV')

add_eyebrow(s, 0.5, 0.7, 'IEU · Pile technique')
add_text(s, 0.5, 1.0, 12, 0.7,
         [('Choix du ', {'font': DISPLAY, 'size': 28, 'color': NAVY, 'bold': True}),
          ('MIE', {'font': DISPLAY, 'size': 28, 'color': YELLOW_DEEP, 'bold': True}),
          (" et récupération de l'INS", {'font': DISPLAY, 'size': 28, 'color': NAVY, 'bold': True})])
add_kicker(s, 0.5, 1.6, 12, 0.55,
           "Deux décisions distinctes. FranceConnect+ ne renvoie pas l'INS. ApCV oui. Stratégie INSi à choisir selon l'IdN.")

# Left: MIE table + numbered récupération
add_text(s, 0.5, 2.3, 6, 0.35,
         [('Choix du MIE patient', {'font': DISPLAY, 'size': 14, 'color': INK, 'bold': True})])

mie_rows = [
    [[('FranceConnect+', {'size': 10, 'bold': True})],
     [('Substantiel', {'size': 10})],
     [('Non', {'size': 10})],
     [('Conv. DINUM + OIDC', {'size': 10})]],
    [[('ApCV', {'size': 10, 'bold': True})],
     [('Substantiel', {'size': 10})],
     [('Oui', {'size': 10, 'color': GREEN, 'bold': True})],
     [('Agrément add. 8 / DI v4.00', {'size': 10})]],
    [[('MdP + OTP', {'size': 10, 'bold': True})],
     [('Faible-Subst.', {'size': 10})],
     [('Non', {'size': 10})],
     [('Maison (PGSSI‑S)', {'size': 10})]],
    [[('eID nationale', {'size': 10, 'bold': True})],
     [('Élevé', {'size': 10})],
     [('Non', {'size': 10})],
     [('Niche', {'size': 10})]],
]
add_table(s, 0.5, 2.7, 6.0, 1.85, mie_rows,
          col_widths=[1.5, 1.3, 0.8, 2.4],
          header=['MIE', 'eIDAS', 'INS ?', 'Effort'])

add_text(s, 0.5, 4.75, 6, 0.35,
         [('Récupération du matricule INS', {'font': DISPLAY, 'size': 14, 'color': INK, 'bold': True})])
numbered_list(s, 0.5, 5.15, 6, 1.85,
              [[("Référentiel d'identités", True), (" : appel INSi via CPx ou IGC‑Santé (INS 37, 39, 40).", False)],
               [("Consommateur d'identité", True), (" : flux IHE PAM en réception (INS 45).", False)],
               [("Auto-saisie traits", True), (" = bascule en référentiel d'identités (Guide INS V3.0 §1.2).", False)]],
              size=10)

# Right: warm block (piège IEU 8) + dark block (2FA composition)
add_block(s, 6.85, 2.3, 6.0, 1.9, 'warm')
block_header(s, 7.1, 2.45, 5.5, 'Piège IEU 8')
add_text(s, 7.1, 2.8, 5.5, 1.3,
         [('FranceConnect+ seul = identité civile ', {'font': BODY, 'size': 11, 'color': INK_BODY}),
          ('sans INS', {'font': BODY, 'size': 11, 'color': INK, 'bold': True}),
          (". IEU 8 oblige à le récupérer ailleurs. Sinon chaînage DMP / Mon espace santé impossible.",
           {'font': BODY, 'size': 11, 'color': INK_BODY})])

add_block(s, 6.85, 4.35, 6.0, 2.55, 'dark')
block_dark_header(s, 7.1, 4.5, 5.5, '2FA conforme : composition')
bullet_list(s, 7.1, 4.85, 5.7, 2.0,
            [[('Facteur 1', True), (' : MdP ', False), ('ou', True), (' smartphone enrôlé.', False)],
             [('Facteur 2', True), (' distinct : TOTP, SMS sur n° vérifié, push signé.', False)],
             [('SMS', True), (' : OK si n° vérifié + révocable (IEU 6). À renforcer.', False)],
             [('Email seul comme 2ᵉ facteur', True), (' : non conforme.', False)]],
            dark=True, size=10)

# ============================================================
# SLIDE 10 · PLAN D'ACTION
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.333, 7.5, PAPER)
add_strip(s, "10 · SUNRISE", "PLAN D'ACTION")
add_foot(s, '10 / 10', "PLAN D'ACTION 8 SEMAINES", "THEODO HEALTHTECH · QARA")

add_eyebrow(s, 0.5, 0.7, "Sunrise · Plan d'action")
add_text(s, 0.5, 1.0, 12, 0.8,
         [("Plan d'action ", {'font': DISPLAY, 'size': 30, 'color': NAVY, 'bold': True}),
          ('8 semaines', {'font': DISPLAY, 'size': 30, 'color': YELLOW_DEEP, 'bold': True})])
add_kicker(s, 0.5, 1.65, 12, 0.55,
           'Six chantiers en parallèle. PSC et INSi sont sur le chemin critique : démarrer en S1.')

# Left: numbered chantiers
add_text(s, 0.5, 2.4, 6, 0.35,
         [('Chantiers', {'font': DISPLAY, 'size': 14, 'color': INK, 'bold': True})])
numbered_list(s, 0.5, 2.8, 6, 4.0,
              [[('Qualifier les profils', True), (' (Accès Pro, Usager, Réf. d\'identités, ApCV).', False)],
               [('Lancer PSC', True), (' sandbox puis prod, rédiger CGU.', False)],
               [('Trancher la stratégie INS', True), (' (INSi vs flux IHE PAM).', False)],
               [('Choisir la 2FA patient', True), (' (FC+ v1, ApCV v2).', False)],
               [("Aligner politique d'auth", True), (' (MdP, MIE, idle).', False)],
               [('Produire les preuves', True), (' Convergence.', False)]],
              size=11)

# Right: dark block (livrables) + warm block (calendar) + note
add_block(s, 6.85, 2.4, 6.0, 2.2, 'dark')
block_dark_header(s, 7.1, 2.55, 5.5, 'Livrables docs')
bullet_list(s, 7.1, 2.9, 5.5, 1.8,
            ['CGU PSC (PSC 5).',
             'Politique gestion MIE (IEPS 5-7, IEU 4-6).',
             'Politique MdP (PGSSI‑S IE ASPP & Usagers).',
             'Procédure attribution / récupération INS.',
             'Plan d\'identitovigilance (Guide INS V3.0).',
             'Campagne de vérif profils PS vs Annuaire (IEPS 4).'],
            dark=True, size=10)

add_block(s, 6.85, 4.75, 6.0, 1.85, 'warm')
block_header(s, 7.1, 4.9, 5.5, 'Calendrier')
cal_rows = [
    [[('S1-S2', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Qualif. produit + iSC / PSC sandbox', {'size': 10})]],
    [[('S3-S4', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Intégration PSC + rédaction politiques', {'size': 10})]],
    [[('S3-S5', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Décision INS + contrat INSi', {'size': 10})]],
    [[('S5-S7', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('2FA patient + récupération INS', {'size': 10})]],
    [[('S7-S8', {'font': MONO, 'size': 10, 'color': NAVY, 'bold': True})],
     [('Preuves + dépôt Convergence', {'size': 10})]],
]
add_table(s, 7.1, 5.25, 5.5, 1.25, cal_rows, col_widths=[1.1, 4.4])

add_text(s, 6.85, 6.75, 6.0, 0.4,
         [('Prochaine étape : atelier 90 min. Livrable : matrice exigences-preuves dans Notion.',
           {'font': BODY, 'size': 10, 'color': INK_SOFT, 'italic': True})])

# ============================================================
# TODO: change to your output path
out = '/Users/nicolasbertrand/Documents/deck-native.pptx'
prs.save(out)
import os
print('written:', out, 'size:', os.path.getsize(out), 'bytes')
